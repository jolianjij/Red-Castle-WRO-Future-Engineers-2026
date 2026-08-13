#!/usr/bin/env python3
"""
build_pptx_journal.py - builds Engineering-Journal.pptx from the repo's own
documentation (README.md, ENGINEERING-JOURNAL.md, config.py, the BOM) and the
available photos/renders, following the section structure used in the team's
2025 documentation (Team -> Introduction -> Solution -> Mobility Management ->
Power and Sense Management -> Robot's Evolution -> Photos -> Links).

Run:  python build_pptx_journal.py
Produces: Engineering-Journal.pptx  (open in PowerPoint, or run
          build_pptx_to_pdf.py / export manually to get the PDF for the
          hard-copy submission).
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.abspath(__file__))

RED = RGBColor(0xA4, 0x1D, 0x24)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF3, 0xF0, 0xEE)
AMBER = RGBColor(0xE4, 0x9B, 0x1F)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
FLOW_DECISION = RGBColor(0xFB, 0xEC, 0xD9)
FLOW_ACTION = RGBColor(0xEF, 0xE3, 0xE1)
FLOW_TERM = RGBColor(0xA4, 0x1D, 0x24)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def img(name):
    for sub in ("v-photos", "t-photos", "models/renders", "models/renders/parts",
                "schemes", "other", "other/testing-photos"):
        p = os.path.join(ROOT, sub, name)
        if os.path.exists(p):
            return p
    return None


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def textbox(slide, l, t, w, h, text, size=18, bold=False, color=DARK,
           align=PP_ALIGN.LEFT, font="Calibri", italic=False, anchor=None):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color
            r.font.name = font
    return tb


def bullets(slide, l, t, w, h, items, size=15, color=DARK, gap=6, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        indent = it.startswith("  ")
        p.text = ("•  " if not indent else "     -  ") + it.strip()
        p.space_after = Pt(gap)
        for r in p.runs:
            r.font.size = Pt(size if not indent else size - 1)
            r.font.color.rgb = color
            r.font.name = font
    return tb


def header_bar(slide, title, subtitle=None, number=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(1.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = RED
    bar.line.fill.background()
    bar.shadow.inherit = False
    textbox(slide, Inches(0.5), Inches(0.10), Inches(11.5), Inches(0.6),
            title, size=26, bold=True, color=WHITE)
    if subtitle:
        textbox(slide, Inches(0.5), Inches(0.62), Inches(11.5), Inches(0.4),
                subtitle, size=13, color=RGBColor(0xF3, 0xD9, 0xDA))
    if number:
        textbox(slide, Inches(12.55), Inches(7.05), Inches(0.6), Inches(0.35),
                str(number), size=11, color=GREY, align=PP_ALIGN.RIGHT)


def picture_fit(slide, path, l, t, w, h, caption=None):
    if not path or not os.path.exists(path):
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        box.fill.solid(); box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = GREY; box.line.width = Pt(1)
        name = caption or "image"
        textbox(slide, l, t + h / 2 - Inches(0.3), w, Inches(0.6),
                f"[ INSERT IMAGE: {name} ]", size=13, italic=True,
                color=GREY, align=PP_ALIGN.CENTER)
        return
    from PIL import Image
    iw, ih = Image.open(path).size
    ar = iw / ih
    box_ar = w / h
    if ar > box_ar:
        pw = w; ph = Emu(int(w / ar))
    else:
        ph = h; pw = Emu(int(h * ar))
    pl = l + (w - pw) // 2
    pt = t + (h - ph) // 2
    slide.shapes.add_picture(path, pl, pt, width=pw, height=ph)
    if caption:
        textbox(slide, l, t + h + Inches(0.05), w, Inches(0.35), caption,
                size=12, italic=True, color=GREY, align=PP_ALIGN.CENTER)


n = [0]
def num():
    n[0] += 1
    return n[0]


# ==========================================================================
# FLOWCHART HELPERS (native pptx shapes, no external renderer needed)
# ==========================================================================
def flow_box(slide, cx, cy, w, h, text, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             fill=FLOW_ACTION, line=RED, size=9.5, textcolor=DARK, bold=False):
    l, t = int(cx - w / 2), int(cy - h / 2)
    sh = slide.shapes.add_shape(shape, l, t, int(w), int(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(3); tf.margin_right = Pt(3)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line_txt in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line_txt
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = textcolor
            r.font.bold = bold
            r.font.name = "Calibri"
    return {"c": (cx, cy), "l": (l, cy), "r": (l + w, cy),
            "t": (cx, t), "b": (cx, t + h), "w": w, "h": h}


def flow_arrow(slide, pts, color=GREY, label=None, dashed=False):
    pts = [(int(x), int(y)) for x, y in pts]
    for i in range(len(pts) - 1):
        x1, y1 = pts[i]
        x2, y2 = pts[i + 1]
        if x1 == x2 and y1 == y2:
            continue
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        conn.line.color.rgb = color
        conn.line.width = Pt(1.25)
        ln = conn.line._get_or_add_ln()
        if dashed:
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
        if i == len(pts) - 2:
            ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    if label:
        mi = (len(pts) - 1) // 2
        mx = (pts[mi][0] + pts[mi + 1][0]) / 2
        my = (pts[mi][1] + pts[mi + 1][1]) / 2
        box = slide.shapes.add_textbox(int(mx - Inches(0.55)), int(my - Inches(0.14)),
                                        Inches(1.1), Inches(0.26))
        tf = box.text_frame; tf.word_wrap = False
        tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
        p = tf.paragraphs[0]; p.text = label; p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = color; r.font.name = "Calibri"
        box.fill.solid(); box.fill.fore_color.rgb = WHITE
        box.line.fill.background()


# ==========================================================================
# COVER
# ==========================================================================
s = add_slide(); bg(s, DARK)
logo = img("team-logo.jpg")
if logo:
    picture_fit(s, logo, Inches(3.67), Inches(0.8), Inches(6), Inches(2.7))
textbox(s, Inches(1), Inches(3.7), Inches(11.33), Inches(0.9),
        "Engineering Journal", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(4.5), Inches(11.33), Inches(0.6),
        "WRO 2026 Future Engineers — Self-Driving Cars", size=20, color=RGBColor(0xE0,0xB8,0xB9),
        align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(5.5), Inches(11.33), Inches(0.5),
        "Team The Red Castle  |  HMK AI and Robotics Club", size=16, color=WHITE,
        align=PP_ALIGN.CENTER, bold=True)
textbox(s, Inches(1), Inches(6.0), Inches(11.33), Inches(0.5),
        "Coach: Ahmad Kalthom   |   Jolian Wassof  ·  Omar Shammout  ·  Louay Rashwan",
        size=14, color=RGBColor(0xCC,0xCC,0xCC), align=PP_ALIGN.CENTER)

# ==========================================================================
# TABLE OF CONTENTS
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "Table of Contents")
toc = [
    "1.  Introduction",
    "2.  Solution",
    "      2.1  Open Challenge   2.2  Obstacle Challenge",
    "      2.3  Control-Loop Flowchart (Sense & Decide / Act & Loop)",
    "      2.4  Obstacle Challenge Priority Flowchart",
    "3.  Mobility Management",
    "      3.1  The Steering Assembly   3.2  The Driving Assembly",
    "      3.3  The Robot's Body / 3D Design   3.3b  Printed Parts Gallery",
    "      3.4  Mobility Measurement",
    "4.  Power and Sense Management",
    "      4.1  The Robot's Power Source   4.2  Main Components",
    "      4.2b  Why We Chose Each Component",
    "      4.3  Wiring Diagram   4.4  Sensing",
    "5.  The Robot's Evolution",
    "      5.1  Past   5.1a  Camera Decision — Proof",
    "      5.1b/c  Vision Testing — Flip, Colour, Masking, Calibration",
    "      5.2  Present   5.3  Future",
    "6.  The Vehicle's Photos",
    "7.  Team",
    "8.  Links and Credits",
]
bullets(s, Inches(0.9), Inches(1.4), Inches(11), Inches(5.9), toc, size=15.5, gap=7)

# ==========================================================================
# 1. INTRODUCTION
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "1. Introduction", number=num())
textbox(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(5.6),
    "Self-driving cars are one of the clearest places where robotics, computer "
    "vision and control theory meet in a single physical object that either "
    "works or doesn't. For WRO 2026 Future Engineers, our goal was to build a "
    "vehicle that finds its own way around a track it has never measured, using "
    "only what a camera can see.\n\n"
    "We built our car around a Raspberry Pi 4 and a single wide-angle camera. "
    "No LiDAR, no ultrasonic rangefinders, no wheel encoders, no IMU. Everything "
    "the rules require the car to know — where the walls are, which way to lap, "
    "where the traffic signs are, where to park — is visible to a camera, so we "
    "put our engineering time into making that one sensor trustworthy instead "
    "of into fusing several imperfect ones.\n\n"
    "That decision shaped most of what follows. A camera-only car lives or dies "
    "on image quality and threshold tuning, and a large part of this document is "
    "the story of getting those right: a sensor swap that fixed a depth-of-field "
    "problem we didn't expect, a wall detector that turned out to be reading "
    "corner lines as walls, and a control law that we rebuilt once we had a "
    "correct measurement to calibrate it against.\n\n"
    "The vehicle drives the Open Challenge — three laps around a track whose "
    "inner wall positions are randomised each round — and the Obstacle Challenge, "
    "which adds red and green traffic signs that must be passed on the correct "
    "side and a parking manoeuvre at the end.",
    size=15.5, color=DARK)

# ==========================================================================
# 2. SOLUTION - OPEN CHALLENGE
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "2. Solution", "2.1  Open Challenge", number=num())
textbox(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(0.9),
    "The Open Challenge is three laps of a track whose inner wall positions are "
    "randomised each round, so nothing about the corridor width can be hard-coded "
    "in advance — the car has to read it from the camera as it drives.",
    size=15, color=DARK)
bullets(s, Inches(0.7), Inches(2.3), Inches(11.9), Inches(4.6), [
    "Direction is decided once, from the first corner line the camera sees, and locked "
    "for the rest of the run — whichever colour reads bigger in that frame sets clockwise "
    "or counter-clockwise, and the reading is cross-checked against which side of the "
    "image has more open wall as a second opinion.",
    "Between corners, a proportional-derivative controller holds a fixed distance to the "
    "outer wall only. We deliberately do not centre between both walls: the inner wall "
    "moves every round, so the outer wall is the one reference that stays put.",
    "At a corner, crossing the driving-direction's line fires a scripted turn — full "
    "steering lock held for a bounded time, ending as soon as the way ahead reads clear "
    "rather than after a fixed duration.",
    "A wall emergency always wins, in any state: if either wall gets too close the car "
    "escapes proportionally to how close it is, and if both walls are close at once "
    "(jammed facing a corner) it commits to a single escape direction instead of "
    "wavering between two.",
    "The car stops once 11 quadrants are counted and the lap timer has run out, so it "
    "finishes driving through the last straight instead of halting on the line itself.",
], size=15, gap=14)

# ==========================================================================
# 2.2 OBSTACLE CHALLENGE
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "2. Solution", "2.2  Obstacle Challenge", number=num())
textbox(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(0.9),
    "The Obstacle Challenge adds red and green traffic signs that must be passed on "
    "the correct side, then a parking manoeuvre once the laps are done.",
    size=15, color=DARK)
bullets(s, Inches(0.7), Inches(2.3), Inches(11.9), Inches(4.6), [
    "If a sign is visible, the car steers to place it correctly — red on the car's right, "
    "green on its left — with the target position sliding further toward the frame edge "
    "as the sign gets closer, so the manoeuvre curves rather than swerving late.",
    "If a sign was visible a moment ago but drops out of view for a frame, the last "
    "steering command is held rather than reacted to — sign detection can flicker sharply "
    "between consecutive frames, and reacting to every drop-out broke passes mid-way.",
    "Between signs, the car falls back to the same outer-wall lane-keeping law used in the "
    "Open Challenge, instead of driving straight and hoping a wall override corrects it in "
    "time.",
    "A wall override sits above all of this and also does the job of turning the car "
    "through each corner, since there is no separate scripted-turn sequencer here.",
    "Parking is not yet wired into the main loop — the magenta gate is currently treated "
    "purely as a wall to avoid. This is called out honestly in the Evolution section.",
], size=15, gap=14)

# ==========================================================================
# 2.3 CONTROL-LOOP FLOWCHART I
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "2. Solution", "2.3  Control-Loop Flowchart — Sense & Decide", number=num())
textbox(s, Inches(0.5), Inches(1.12), Inches(12.3), Inches(0.3),
    "Every frame runs the same read -> decide -> act loop. navigate() in robot.py picks the "
    "steering command in this strict priority order.", size=12.5, italic=True, color=GREY)

LX = Inches(2.55)
RX = Inches(9.9)
A = flow_box(s, LX, Inches(1.62), Inches(2.1), Inches(0.34), "Start", shape=MSO_SHAPE.OVAL,
             fill=FLOW_TERM, line=FLOW_TERM, textcolor=WHITE, size=10, bold=True)
B = flow_box(s, LX, Inches(2.18), Inches(3.1), Inches(0.46), "Capture frame\n(OV5647, 180° flip)")
C = flow_box(s, LX, Inches(2.76), Inches(3.1), Inches(0.46), "Crop mat ROI -> 320x160,\nmedian blur, convert HSV")
D = flow_box(s, LX, Inches(3.34), Inches(3.1), Inches(0.46), "Measure L/R wall density,\nfront density, line fractions")
E = flow_box(s, LX, Inches(3.92), Inches(3.1), Inches(0.46), "LapTracker.update:\ndirection + lap counting")
F = flow_box(s, LX, Inches(4.62), Inches(3.1), Inches(0.62), "Either wall past\nWALL_EMERGENCY?",
             shape=MSO_SHAPE.DIAMOND, fill=FLOW_DECISION)
H = flow_box(s, LX, Inches(5.34), Inches(3.1), Inches(0.6), "Scripted turn\nactive?",
             shape=MSO_SHAPE.DIAMOND, fill=FLOW_DECISION)
J = flow_box(s, LX, Inches(6.04), Inches(3.1), Inches(0.6), "Front wall very close\n& direction known?",
             shape=MSO_SHAPE.DIAMOND, fill=FLOW_DECISION)
L = flow_box(s, LX, Inches(6.68), Inches(3.1), Inches(0.42), "Outer-wall PD control:\nhold distance to ONE wall",
             fill=RGBColor(0xE3, 0xEE, 0xE2))

G = flow_box(s, RX, Inches(4.62), Inches(2.9), Inches(0.5), "Proportional escape;\nfull lock if BOTH walls close")
I = flow_box(s, RX, Inches(5.34), Inches(2.9), Inches(0.5), "Hold full lock in the\nlocked direction")
K = flow_box(s, RX, Inches(6.04), Inches(2.9), Inches(0.5), "Trigger the turn anyway\n(geometry backup)")

flow_arrow(s, [A["b"], B["t"]])
flow_arrow(s, [B["b"], C["t"]])
flow_arrow(s, [C["b"], D["t"]])
flow_arrow(s, [D["b"], E["t"]])
flow_arrow(s, [E["b"], F["t"]])
flow_arrow(s, [F["b"], H["t"]], label="no")
flow_arrow(s, [H["b"], J["t"]], label="no")
flow_arrow(s, [J["b"], L["t"]], label="no")
flow_arrow(s, [F["r"], G["l"]], label="yes")
flow_arrow(s, [H["r"], I["l"]], label="yes")
flow_arrow(s, [J["r"], K["l"]], label="yes")

note = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(7.12), Inches(10.6), Inches(0.3))
note.fill.solid(); note.fill.fore_color.rgb = LIGHT
note.line.color.rgb = RED; note.line.width = Pt(1)
ln = note.line._get_or_add_ln(); ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
note.shadow.inherit = False
tf = note.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.text = "G, I, K or L → the steering command continues on “Control-Loop Flowchart — Act & Loop”"
p.alignment = PP_ALIGN.CENTER
for r in p.runs:
    r.font.size = Pt(11); r.font.italic = True; r.font.color.rgb = RED; r.font.name = "Calibri"

# ==========================================================================
# 2.3b CONTROL-LOOP FLOWCHART II
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "2. Solution", "2.3  Control-Loop Flowchart — Act & Loop", number=num())
textbox(s, Inches(0.5), Inches(1.12), Inches(12.3), Inches(0.3),
    "Continues from the steering command chosen on the previous slide (G, I, K or L).",
    size=12.5, italic=True, color=GREY)

CX = Inches(6.6)
M0 = flow_box(s, CX, Inches(1.75), Inches(3.4), Inches(0.4), "steering command", shape=MSO_SHAPE.OVAL,
              fill=WHITE, line=GREY, textcolor=GREY, size=10)
M = flow_box(s, CX, Inches(2.4), Inches(4.2), Inches(0.46), "Clamp to STEER_MAX, set servo")
N = flow_box(s, CX, Inches(3.0), Inches(4.2), Inches(0.46), "Scale speed down with steering\neffort, set motor")
O = flow_box(s, CX, Inches(3.6), Inches(4.2), Inches(0.46), "Log the frame to CSV")
P = flow_box(s, CX, Inches(4.35), Inches(4.2), Inches(0.65), "Lap count satisfied\nand timer expired?",
             shape=MSO_SHAPE.DIAMOND, fill=FLOW_DECISION)
Q = flow_box(s, CX, Inches(5.15), Inches(3.0), Inches(0.42), "Stop: motor 0, servo centre",
             shape=MSO_SHAPE.OVAL, fill=FLOW_TERM, line=FLOW_TERM, textcolor=WHITE, size=10.5, bold=True)

flow_arrow(s, [M0["b"], M["t"]])
flow_arrow(s, [M["b"], N["t"]])
flow_arrow(s, [N["b"], O["t"]])
flow_arrow(s, [O["b"], P["t"]])
flow_arrow(s, [P["b"], Q["t"]], label="yes")
loopX = Inches(10.6)
flow_arrow(s, [P["r"], (loopX, P["c"][1]), (loopX, M0["c"][1]), M0["r"]], label="no — next frame", dashed=True)

textbox(s, Inches(1.2), Inches(5.9), Inches(10.9), Inches(1.2),
    "This priority order is what makes a wall emergency unconditional: it is checked first, "
    "every frame, before the turn state or the lane-keeping law ever gets a say. The scripted "
    "turn and the geometry backup only run when no wall is already too close.",
    size=13.5, color=DARK, italic=True)

# ==========================================================================
# 2.4 OBSTACLE CHALLENGE PRIORITY FLOWCHART
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "2. Solution", "2.4  Obstacle Challenge — Priority Flowchart", number=num())
textbox(s, Inches(0.5), Inches(1.12), Inches(12.3), Inches(0.3),
    "Steering is decided by a strict priority, so a lower-priority behaviour can never "
    "override safety.", size=12.5, italic=True, color=GREY)

OCX = Inches(6.6)
oA = flow_box(s, OCX, Inches(1.75), Inches(2.6), Inches(0.36), "Every frame", shape=MSO_SHAPE.OVAL,
              fill=FLOW_TERM, line=FLOW_TERM, textcolor=WHITE, size=10, bold=True)
oB = flow_box(s, OCX, Inches(2.35), Inches(3.4), Inches(0.55), "Sign visible?",
              shape=MSO_SHAPE.DIAMOND, fill=FLOW_DECISION)

cC = flow_box(s, Inches(2.6), Inches(3.35), Inches(3.5), Inches(0.6),
              "Steer to place it correctly:\nred on our right, green on our left")
cD = flow_box(s, Inches(6.6), Inches(3.35), Inches(3.5), Inches(0.6),
              "No, but seen recently:\nhold last manoeuvre (SIGN_MEMORY)")
cE = flow_box(s, Inches(10.6), Inches(3.35), Inches(3.4), Inches(0.6),
              "No sign at all:\nfall back to outer-wall lane keeping")

oF = flow_box(s, OCX, Inches(4.4), Inches(3.6), Inches(0.6), "Outer wall past\nWALL_EMERGENCY?",
              shape=MSO_SHAPE.DIAMOND, fill=FLOW_DECISION)
oG = flow_box(s, Inches(4.3), Inches(5.35), Inches(4.6), Inches(0.55),
              "Full lock away from the wall —\noverrides everything, also turns corners",
              fill=RGBColor(0xF7, 0xE0, 0xDF))
oH = flow_box(s, Inches(9.3), Inches(5.35), Inches(4.2), Inches(0.55), "Use the steering chosen above",
              fill=RGBColor(0xE3, 0xEE, 0xE2))

flow_arrow(s, [oA["b"], oB["t"]])
flow_arrow(s, [oB["l"], (Inches(2.6), oB["c"][1]), cC["t"]], label="yes")
flow_arrow(s, [oB["b"], cD["t"]], label="no, recent")
flow_arrow(s, [oB["r"], (Inches(10.6), oB["c"][1]), cE["t"]], label="no")
flow_arrow(s, [cC["b"], (cC["c"][0], Inches(3.9)), (OCX, Inches(3.9)), oF["t"]])
flow_arrow(s, [cD["b"], oF["t"]])
flow_arrow(s, [cE["b"], (cE["c"][0], Inches(3.9)), (OCX, Inches(3.9)), oF["t"]])
flow_arrow(s, [oF["l"], (Inches(4.3), oF["c"][1]), oG["t"]], label="yes")
flow_arrow(s, [oF["r"], (Inches(9.3), oF["c"][1]), oH["t"]], label="no")

textbox(s, Inches(1.0), Inches(6.15), Inches(11.3), Inches(1.0),
    "Priority, highest to lowest: (1) wall override — also turns every corner, "
    "(2) sign steering, (3) sign-hold memory for detection flicker, (4) outer-wall lane "
    "keeping between signs.", size=13, color=GREY, italic=True)

# ==========================================================================
# 3.1 STEERING ASSEMBLY
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "3. Mobility Management", "3.1  The Steering Assembly", number=num())
textbox(s, Inches(0.6), Inches(1.3), Inches(6.6), Inches(5.7),
    "Steering is Ackermann-style, driven by a single SG90 micro servo (GPIO13, "
    "which is a hardware-PWM pin — a smoother, lower-jitter signal than a "
    "software-timed one).\n\n"
    "The servo turns a central bell-crank, which pushes a tie-rod out to both "
    "front steering knuckles. Because the linkage geometry is Ackermann and not "
    "a simple parallel arm, the inner wheel turns through a sharper angle than "
    "the outer one in a corner — each wheel follows its own turning radius "
    "instead of fighting the other and scrubbing.\n\n"
    "The linkage's mechanical limit is ±35°. In software we run it at ±20° "
    "(STEER_MAX in config.py) for stability at full driving speed — the wider "
    "mechanical range is there if a future tuning pass wants it.\n\n"
    "Every steering command in the codebase passes through this one constant and "
    "one servo() function, so there is exactly one place that ever needs to know "
    "the true centre position — the field-calibrated trim is stored in "
    "servo_center.txt (currently -9°) and applied automatically.",
    size=14.5, color=DARK)
picture_fit(s, img("front-steering-linkage.png"), Inches(7.5), Inches(1.3),
           Inches(5.2), Inches(2.55), "Front steering linkage, top view")
picture_fit(s, img("chassis-top-steering-left.png"), Inches(7.5), Inches(4.35),
           Inches(5.2), Inches(2.55), "Steering turned, showing the linkage sweep")

# ==========================================================================
# 3.2 DRIVING ASSEMBLY
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "3. Mobility Management", "3.2  The Driving Assembly", number=num())
textbox(s, Inches(0.6), Inches(1.3), Inches(6.6), Inches(5.7),
    "Drive comes from an N20 geared DC motor rated 12 V, 200 rpm, chosen for its "
    "compact size and enough torque to pull the car away from rest repeatedly "
    "without stalling.\n\n"
    "The motor is controlled by an L9110S dual H-bridge, wired so one input is "
    "held low while the other carries the PWM speed signal (GPIO24 = forward, "
    "GPIO23 = reverse). In software, motor() never switches directly from "
    "forward to reverse — it coasts to a stop and waits first, because slamming "
    "the direction while the motor is still spinning drives its back-EMF onto "
    "the supply rail and can destroy the Pi's regulator. We found this out by "
    "losing a regulator early in the build.\n\n"
    "From the motor, a 25-tooth pinion drives a 20-tooth gear (a 1.25:1 "
    "reduction) into a 3D-printed differential we designed ourselves — outer "
    "shell, outer ring gear, and two inner bevel gears on long and short shafts. "
    "The differential is the single change that let us raise the steering limit: "
    "with a solid axle, the outer wheel scrubbing against the inner one through "
    "a corner cost so much traction we had to cap steering near 8° just to keep "
    "the car from stalling mid-turn.",
    size=14, color=DARK)
picture_fit(s, img("rear-axle-differential-mesh.png"), Inches(7.5), Inches(1.3),
           Inches(5.2), Inches(2.55), "N20 motor driving the differential gear mesh")
picture_fit(s, img("rear-axle-differential-front.png"), Inches(7.5), Inches(4.35),
           Inches(5.2), Inches(2.55), "Rear axle, front view")

# ==========================================================================
# 3.3 ROBOT'S BODY / 3D DESIGN
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "3. Mobility Management", "3.3  The Robot's Body / 3D Design", number=num())
textbox(s, Inches(0.6), Inches(1.3), Inches(11.9), Inches(0.85),
    "The chassis was designed in Fusion 360 and printed in PLA+ Silk Silver on a "
    "Bambu Lab A1: a base plate, a middle plate and a top plate carrying the Pi, "
    "battery, driver, buck converter and the camera mast, sized to keep mass low "
    "and central over the wheelbase.",
    size=14.5, color=DARK)
picture_fit(s, img("full-assembly.png"), Inches(1.4), Inches(2.2),
           Inches(10.5), Inches(2.85), "Full car assembly")
textbox(s, Inches(0.6), Inches(5.65), Inches(11.9), Inches(1.5),
    "Mass 407 g, overall footprint 14.2 x 9.3 x 15 cm, wheelbase 9.4 cm, rear "
    "track 8.5 cm — a wheelbase-to-track ratio close to 1:1, a stable "
    "proportion for cornering without tipping risk at this size.",
    size=13.5, color=GREY, italic=True)

# ==========================================================================
# 3.3b PRINTED PARTS GALLERY
# ==========================================================================
part_photos = [
    ("base.png", "Base"), ("middle-plate.png", "Middle plate"),
    ("top-plate.png", "Top plate"), ("camera-holder.png", "Camera holder"),
    ("battery-slider.png", "Battery slider"), ("n20-motor-holder.png", "N20 motor holder"),
    ("n20-gear.png", "N20 gear"), ("inner-differential-gear.png", "Inner differential gear"),
    ("outer-differential-gear.png", "Outer differential gear"),
]
part_photos_2 = [
    ("differential-outer-shell.png", "Differential outer shell"),
    ("inner-gear-long-shaft.png", "Inner gear, long shaft"),
    ("inner-gear-short-shaft.png", "Inner gear, short shaft"),
    ("ring.png", "Ring"), ("steering-base.png", "Steering base"),
    ("steering-base-servo-mount.png", "Steering base servo mount"),
    ("steering-nut.png", "Steering nut"), ("steering-wheel-mount.png", "Steering wheel mount"),
    ("wheel.png", "Wheel"),
]
for page_num, plist in enumerate((part_photos, part_photos_2), start=1):
    s = add_slide(); bg(s)
    header_bar(s, "3. Mobility Management",
               "3.3b  Printed Parts — Individual Renders (%d/2)" % page_num, number=num())
    cols = 3
    pw, ph = Inches(3.9), Inches(1.5)
    gx = Inches(0.25)
    row_pitch = Inches(2.0)
    x0, y0 = Inches(0.7), Inches(1.2)
    for i, (fname, label) in enumerate(plist):
        c, r = i % cols, i // cols
        l = x0 + c * (pw + gx)
        t = y0 + r * row_pitch
        picture_fit(s, img(fname), l, t, pw, ph, label)

# ==========================================================================
# 3.4 MOBILITY MEASUREMENT
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "3. Mobility Management", "3.4  Mobility Measurement", number=num())
textbox(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(0.6),
    "To reason about how fast the car can safely go, we worked through the "
    "drivetrain's speed chain from the motor to the ground:",
    size=15, color=DARK)
textbox(s, Inches(1.0), Inches(2.15), Inches(10.8), Inches(4.6),
    "Wheel radius (rw):\n"
    "      rw = 4.7 cm / 2 = 2.35 cm = 0.0235 m\n\n"
    "Wheel perimeter (Pw):\n"
    "      Pw = 2π·rw = 2 × 3.1416 × 0.0235 ≈ 0.1477 m\n\n"
    "Motor speed and the gear reduction:\n"
    "      ωmotor = 200 rpm  (rated, no load)\n"
    "      ωaxle = ωmotor / 1.25  (25:20 reduction)  = 160 rpm\n\n"
    "Ground speed (theoretical, no load):\n"
    "      v = Pw · ωaxle = 0.1477 × 160 ≈ 23.6 m·min⁻¹\n"
    "      v ≈ 0.394 m·s⁻¹  (≈ 39 cm/s)\n\n"
    "This is a ceiling, not a running speed: tyre friction, the differential, "
    "and the car's own 407 g mass all take a share of it. What actually governs "
    "how fast the car drives in competition is the 100% cruise setting in "
    "software, which cruise_speed() then eases back automatically whenever the "
    "steering is working hard — fast on the straights, deliberately slower "
    "through a turn.",
    size=15, color=DARK, font="Consolas")

# ==========================================================================
# 4.1 POWER SOURCE
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "4. Power and Sense Management", "4.1  The Robot's Power Source", number=num())
textbox(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(5.7),
    "Power comes from three 18650 Li-ion cells in series — about 11.1 V nominal, "
    "up to 12.6 V fully charged — feeding two separate rails rather than one "
    "shared one:\n\n"
    "The motor is wired straight from the battery to the L9110S driver's supply "
    "pin, with nothing in between. A separate DC-DC step-down converter, rated "
    "at 5 V and at least 3 A, powers the Raspberry Pi and the servo on its own "
    "rail.\n\n"
    "This split exists because of a real failure: early on, the motor and the "
    "Pi shared one supply, and the Pi rebooted every time the motor drew current "
    "under load — a classic brownout, made worse by the driver's ground not "
    "being tied back to the Pi's ground at all, so the only electrical link "
    "between them was the GPIO signal wire itself. Separating the rails and "
    "tying every ground together (battery, driver, Pi) fixed it completely, and "
    "is now a fixed rule on the car rather than a tuning choice.\n\n"
    "A rocker switch sits in the battery line as the master on/off. We also note, "
    "as a live watch item rather than a solved problem: a freshly charged pack "
    "at 12.6 V sits marginally above the L9110S's 12 V rating, so driver "
    "temperature is worth checking after a run.",
    size=15, color=DARK)

# ==========================================================================
# 4.2 MAIN COMPONENTS
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "4. Power and Sense Management", "4.2  Main Components", number=num())
bullets(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(5.7), [
    "OV5647 wide-angle camera — a small-sensor, fixed-focus module with roughly a "
    "120° field of view. It is the vehicle's only sensor: every wall, corner line, "
    "traffic sign and the parking gate is read from its image. Fixed focus means a "
    "large depth of field, so the near mat and a far wall stay sharp at the same "
    "time with nothing to hunt or mis-lock mid-run — a deliberate choice explained "
    "in the Evolution section.",
    "Raspberry Pi 4 Model B, 2 GB RAM, running Raspberry Pi OS Bookworm — the single "
    "board that captures every frame, runs the vision pipeline, decides the steering "
    "and speed, and drives the servo and motor through its GPIO header.",
    "L9110S dual H-bridge motor driver — a compact two-pin-per-motor interface, "
    "chosen for its small size and simple PWM control, sitting on its own supply "
    "rail straight from the battery.",
    "SG90 micro servo — 9 g, powered from the Pi's 5 V rail, driving the Ackermann "
    "steering linkage.",
], size=15.5, gap=16)

# ==========================================================================
# 4.2b WHY WE CHOSE EACH COMPONENT
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "4. Power and Sense Management", "4.2b  Why We Chose Each Component", number=num())

rationale_rows = [
    ("Camera", "OV5647 fixed-focus, wide FOV",
     "Pi Camera Module 3 (autofocus)",
     "Autofocus hunted and lost distant signs to blur (proof: 5.1a). A ground "
     "robot never needs to refocus, so fixed focus cost us nothing and bought "
     "a depth of field wide enough for the near mat and a far sign at once."),
    ("Motor", "N20, 12 V, 200 rpm geared DC",
     "Same footprint, higher-rpm N20",
     "A faster variant trades torque for speed; we needed the torque to pull "
     "the 407 g car from rest repeatedly without stalling more than we needed "
     "extra top speed."),
    ("Driver", "L9110S dual H-bridge",
     "L298N dual H-bridge",
     "L298N's extra current headroom is wasted on a single N20 and it takes "
     "far more board space; L9110S needs only two GPIO pins and fits the "
     "chassis footprint we had."),
    ("Servo", "SG90 micro servo, GPIO13 (HW-PWM)",
     "MG90S metal-gear servo",
     "The steering linkage's load never approached the point where MG90S's "
     "extra torque or durability would matter; SG90 is lighter and cheap "
     "enough to keep spares on hand."),
    ("Battery", "3 x 18650 Li-ion, series (~11.1 V)",
     "2S/3S LiPo pack",
     "18650 cells swap individually without a balance charger, and 11.1 V "
     "nominal clears the buck converter's dropout with margin without "
     "oversizing the pack."),
    ("Differential", "Custom 3D-printed bevel differential",
     "Solid axle (our first build) / off-the-shelf RC diff",
     "A solid axle scrubbed the outer wheel through every corner and capped "
     "steering near 8°. An off-the-shelf diff didn't match the 25:20 gear "
     "pitch already tooled for the N20, so we designed our own."),
]
tbl_l, tbl_t = Inches(0.5), Inches(1.25)
tbl_w, tbl_h = Inches(12.35), Inches(5.85)
gtbl = s.shapes.add_table(len(rationale_rows) + 1, 4, tbl_l, tbl_t, tbl_w, tbl_h).table
col_w = [Inches(1.55), Inches(3.15), Inches(2.95), Inches(4.7)]
for i, w in enumerate(col_w):
    gtbl.columns[i].width = w
headers = ["Component", "What we chose", "Alternative considered", "Why"]
for i, htext in enumerate(headers):
    cell = gtbl.cell(0, i)
    cell.text = htext
    cell.fill.solid(); cell.fill.fore_color.rgb = RED
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
for ri, row in enumerate(rationale_rows, start=1):
    for ci, text in enumerate(row):
        cell = gtbl.cell(ri, ci)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT if ri % 2 == 0 else WHITE
        cell.margin_left = Pt(5); cell.margin_right = Pt(5)
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for r in p.runs:
                r.font.size = Pt(10.5 if ci == 3 else 11)
                r.font.bold = (ci == 0)
                r.font.color.rgb = DARK
                r.font.name = "Calibri"

# ==========================================================================
# 4.3 WIRING DIAGRAM
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "4. Power and Sense Management", "4.3  Wiring Diagram", number=num())
picture_fit(s, img("wiring-diagram.png"), Inches(1.0), Inches(1.25),
           Inches(11.3), Inches(5.6))

# ==========================================================================
# 4.4 SENSING - VISION PIPELINE
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "4. Power and Sense Management", "4.4  Sensing — Vision Pipeline", number=num())
bullets(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(5.7), [
    "Capture: the camera is read at 640x480 in RGB888, then the frame is "
    "rotated 180° in software, since the camera is physically mounted upside-down "
    "on its mast.",
    "Crop and resize: only the lower portion of the frame — the mat — is kept and "
    "resized down to 320x160 before any colour work happens, which both discards "
    "background clutter above the walls and keeps every frame cheap to process.",
    "Median blur, then HSV: a light blur is applied first because the camera runs "
    "a short exposure and elevated gain to freeze motion, which makes the raw image "
    "noisy enough to break both the colour masks and the wall scan if left "
    "untreated. The frame is then converted to HSV, which is far more stable than "
    "RGB under the mat's own lighting variation.",
    "Wall detection is a two-tier test, not one brightness threshold: a pixel is a "
    "wall if it is very dark regardless of colour, OR dark and low-saturation. A "
    "single brightness cut also lit up the blue and orange corner lines and the "
    "mat's printed markings, since they are dark too — and unevenly between the "
    "two frame halves, which injected a steering bias that looked like a control "
    "problem for a long time before we traced it to the measurement itself.",
    "Colour lines are searched only in the bottom of the region of interest, "
    "because a line painted on the mat physically cannot appear higher in the "
    "frame than the mat itself — searching the whole frame let bluish background "
    "pixels near the top trigger false line readings.",
], size=14.5, gap=11)

# ==========================================================================
# 4.4b SENSING - CAMERA POSITION
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "4. Power and Sense Management", "4.4  Sensing — Camera Position", number=num())
textbox(s, Inches(0.7), Inches(1.3), Inches(11.9), Inches(0.55),
    "CAMERA POSITION IS CRUCIAL FOR THE ROBOT'S PERFORMANCE.", size=17, bold=True, color=RED)
bullets(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(5.0), [
    "Height 12.5 cm above the mat. The walls are 10 cm tall, so the lens must sit "
    "above the wall line to look down onto the mat and see the wall's base — the "
    "white-mat/black-wall edge, the highest-contrast feature in the frame.",
    "Tilt roughly 15° downward, enough to keep the mat and both wall bases in the "
    "lower frame while still seeing far enough ahead to catch a corner line or a "
    "traffic sign before reaching it.",
    "Centred laterally and facing dead ahead, zero yaw. Wall-following compares the "
    "left half of the frame against the right half, so any sideways offset or twist "
    "becomes a permanent steering bias that no amount of threshold tuning removes.",
    "Mounted on a rigid mast, not a flexible arm. A vibrating camera blurs the wall "
    "edge and injects noise straight into the steering loop, so mast stiffness is a "
    "control requirement, not a cosmetic one.",
    "The connection itself uses the Pi's dedicated CSI ribbon port, not a GPIO pin, "
    "so it needs no separate power and does not appear on the wiring diagram.",
], size=14.5, gap=13)

# ==========================================================================
# 5.1 EVOLUTION - PAST
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "5. The Robot's Evolution", "5.1  Past", number=num())
textbox(s, Inches(0.7), Inches(1.3), Inches(11.9), Inches(5.9),
    "Two decisions from earlier in the build turned out to be wrong in ways "
    "that only showed up once we started driving, not while designing.\n\n"
    "The first camera was a Raspberry Pi Camera Module 3 — a larger sensor with "
    "autofocus. Focused on the mat it saw the near track sharply, but distant "
    "traffic signs blurred, and blur desaturates colour badly enough that the "
    "HSV masks lost far-away signs entirely; the car only 'saw' an obstacle once "
    "it was already close. A focus sweep confirmed we genuinely could not have "
    "both near and far sharp on that sensor at once. We replaced it with the "
    "OV5647, a smaller, fixed-focus sensor with a much larger depth of field — "
    "losing autofocus cost nothing, since a ground robot never needs to refocus.\n\n"
    "The second was the drivetrain. Our first rear axle was solid, so both "
    "driven wheels were forced to the same speed. In a corner the outer wheel "
    "has to cover more distance than the inner one, and with no way to do that "
    "it scrubbed, losing traction and sometimes stalling the car mid-turn. Our "
    "first fix was purely in software — cutting the steering limit down step by "
    "step until the scrubbing stopped, which landed near 8°. That kept the "
    "wheels turning but left the car unable to corner tightly enough to follow "
    "the track. It was a workaround for a mechanical problem, not a fix for it.",
    size=14.5, color=DARK)

# ==========================================================================
# 5.1a CAMERA DECISION - VISUAL PROOF
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "5. The Robot's Evolution", "5.1a  Camera Decision — Proof, Not Just a Claim", number=num())
textbox(s, Inches(0.6), Inches(1.12), Inches(12.1), Inches(0.35),
    "Same track, same lighting, same lens height — only the camera module changed. Both "
    "frames are straight off the Pi, unedited.", size=13, italic=True, color=GREY)
picture_fit(s, img("before-module3-blurry.jpg"), Inches(0.6), Inches(1.5),
           Inches(5.95), Inches(3.85), "BEFORE — Camera Module 3 (autofocus), focused on the near mat")
picture_fit(s, img("after-ov5647-sharp.jpg"), Inches(6.75), Inches(1.5),
           Inches(5.95), Inches(3.85), "AFTER — OV5647 (fixed focus), same distance and framing")
textbox(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(1.35),
    "The Module 3's shallow depth of field could hold either the near mat or a distant "
    "sign in focus, never both — a focus sweep across f/2.0-equivalent apertures confirmed "
    "it wasn't a tuning problem. Blur desaturates colour enough that the HSV sign masks "
    "lost far signs entirely. The OV5647's smaller sensor and fixed focus trade autofocus "
    "(never needed by a ground robot) for a depth of field wide enough to keep the whole "
    "corridor sharp at once.", size=13, color=DARK)

# ==========================================================================
# 5.1b VISION TESTING PROCESS
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "5. The Robot's Evolution", "5.1b  Vision Testing — Flip & Colour Tuning", number=num())
picture_fit(s, img("raw-before-180-flip.jpg"), Inches(0.5), Inches(1.15),
           Inches(4.0), Inches(2.3), "Raw sensor — mounted upside-down")
picture_fit(s, img("after-180-flip.jpg"), Inches(4.65), Inches(1.15),
           Inches(4.0), Inches(2.3), "After the 180° software flip")
picture_fit(s, img("color-tuning-v-channel.png"), Inches(8.8), Inches(1.15),
           Inches(4.0), Inches(2.3), "HSV Value-channel scan, tuning")
textbox(s, Inches(0.5), Inches(3.85), Inches(12.3), Inches(0.55),
    "The camera mounts physically upside-down on its mast, so robot.py rotates every "
    "frame 180° before anything else touches it — real capture on the left, corrected on the right.",
    size=12.5, color=DARK)
picture_fit(s, img("color-tuning-green-sign-mask.png"), Inches(0.5), Inches(4.45),
           Inches(4.0), Inches(2.5), "colors.json green-sign mask, live tuning")
textbox(s, Inches(4.75), Inches(4.45), Inches(8.05), Inches(2.75),
    "Colour thresholds are not guessed once and left alone — colour_tuner.py runs against "
    "the live camera so every HSV band (blue line, orange line, red sign, green sign, "
    "magenta gate) is set against the real mat under the real, locked exposure. Each "
    "iteration is checked against a visible mask overlay like the one on the left, so a "
    "threshold that looks right in isolation but bleeds onto the mat itself is caught "
    "before it reaches the car.\n\n"
    "This is the same loop that caught the sign-detection bug in the Present section: the "
    "mat/wall boundary fringe read as green at S 96-111, merging with the real sign (S "
    "135+) into one mis-shaped blob — visible only once we started saving the mask, not "
    "from the numbers alone.",
    size=12.5, color=DARK)

# ==========================================================================
# 5.1c VISION TESTING - MASKING AND CALCULATION
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "5. The Robot's Evolution", "5.1c  Vision Testing — Wall Masking & Distance Calibration", number=num())
picture_fit(s, img("wall-line-mask-overlay.png"), Inches(0.5), Inches(1.15),
           Inches(4.55), Inches(4.0), "shadow_check.py — red = wall, green = free mat")
picture_fit(s, img("wall-boundary-calculation.png"), Inches(5.25), Inches(1.15),
           Inches(3.55), Inches(4.0), "freespace_test.py — gap, steer, boundary")
picture_fit(s, img("distance-calibration-40cm.png"), Inches(9.0), Inches(1.15),
           Inches(3.75), Inches(1.85), "Calibration at a taped 40 cm")
picture_fit(s, img("distance-calibration-25cm.png"), Inches(9.0), Inches(3.45),
           Inches(3.75), Inches(1.85), "Calibration at a taped 25 cm")
textbox(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(1.45),
    "Every distance constant in config.py traces back to frames like these: the car is "
    "placed at a taped, ruler-measured distance from the wall and freespace_test.py logs "
    "the wall-density reading at that exact spot — 40 cm measured 0.1032, 25 cm measured "
    "0.1783 — so OUTER_TARGET is a real calibration curve, not a guessed constant. The "
    "left-hand overlay is the same diagnostic used to catch the wall-detector counting the "
    "blue/orange corner lines and the mat's own printed marks as wall (§9 of the "
    "engineering journal) — red pixels are what the code currently calls \"wall\".",
    size=12.5, color=DARK)

# ==========================================================================
# 5.2 EVOLUTION - PRESENT
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "5. The Robot's Evolution", "5.2  Present", number=num())
textbox(s, Inches(0.7), Inches(1.3), Inches(11.9), Inches(5.9),
    "The current car fixes both of those at the source rather than around them.\n\n"
    "A 3D-printed differential now sits on the rear axle, driven from the N20 "
    "motor through a 25:20 gear reduction. With the wheels free to turn at "
    "different speeds through a corner, the scrub is gone, and the steering "
    "limit came back up from 8° toward the linkage's mechanical ceiling.\n\n"
    "The bigger present-day story, though, was perception rather than "
    "mechanics. Even after the camera swap, the car kept crashing near corners "
    "in both driving directions, in a way that no amount of gain or setpoint "
    "tuning changed. We eventually stripped the controller down to nothing and "
    "started saving an annotated frame every time the steering swung hard, "
    "colouring in every pixel the code was calling 'wall'. The answer was "
    "immediate: the wall detector was any dark pixel, and the blue line, the "
    "orange line, and the mat's own printed markings are all dark. They were "
    "being counted as wall — unevenly between the left and right halves of the "
    "frame — which is exactly why the fault only showed up near corners and "
    "resisted every attempt to tune it away. The controller had never been the "
    "problem; its input was lying to it.\n\n"
    "Once the wall test was corrected to require genuine darkness or darkness "
    "plus low saturation, every distance constant in the project was "
    "re-measured from real centimetres, and the Open Challenge now completes "
    "in both directions with stable control.",
    size=14, color=DARK)

# ==========================================================================
# 5.3 EVOLUTION - FUTURE
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "5. The Robot's Evolution", "5.3  Future", number=num())
bullets(s, Inches(0.7), Inches(1.35), Inches(11.9), Inches(5.6), [
    "Wire the parking manoeuvre into the Obstacle Challenge's main loop. The "
    "magenta gate is detected and currently treated purely as a wall to steer "
    "around; the reserved constants for the approach and stop condition already "
    "exist in config.py, but the behaviour itself is not yet connected.",
    "Close out a shadow-misreading issue seen at two of the four corners in field "
    "testing, where shadow on the mat is occasionally read as part of the wall. "
    "The existing two-tier wall test already rejects most shadow; a dedicated "
    "diagnostic tool (shadow_check.py) was built to separate a genuine wall from "
    "shadow in the field, and tightening the test against this case is active "
    "work.",
    "Move wall-following from a pure proportional-derivative law on wall distance "
    "to a heading-plus-offset controller, so the car corrects its angle to the "
    "wall rather than only its distance from it — this should make the straights "
    "noticeably straighter.",
    "Track a sign's position across several frames instead of frame by frame, to "
    "smooth the approach further than the current short-memory hold already does.",
], size=15.5, gap=17)

# ==========================================================================
# 6. VEHICLE PHOTOS
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "6. The Vehicle's Photos", number=num())
photos = [("front.jpg", "Front"), ("back.jpg", "Back"), ("left.jpg", "Left"),
         ("right.jpg", "Right"), ("top.jpg", "Top"), ("bottom.jpg", "Bottom")]
cols, rows = 3, 2
pw, ph = Inches(3.85), Inches(2.55)
gx, gy = Inches(0.25), Inches(0.25)
x0, y0 = Inches(0.75), Inches(1.35)
for i, (fname, label) in enumerate(photos):
    c, r = i % cols, i // cols
    l = x0 + c * (pw + gx)
    t = y0 + r * (ph + gy + Inches(0.25))
    picture_fit(s, img(fname), l, t, pw, ph - Inches(0.3), label)

# ==========================================================================
# 7. TEAM
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "7. Team", number=num())
members = [
    ("ahmad-kalthom.jpg", "Ahmad Kalthom", "Coach",
     "Computer and Automation Engineering student, Damascus University."),
    ("jolian-wassof.jpg", "Jolian Wassof", "Team member",
     "Al-Awael School. AI, robotics and drones; WRO, Pacific and Syrian AI Olympiads."),
    ("omar-shammout.jpg", "Omar Shammout", "Team member",
     "IT Engineering student, Damascus University. Robotics, AI, electronics."),
    ("louay-rashwan.jpg", "Louay Rashwan", "Team member",
     "Computer and Automation Engineering student, Damascus University."),
]
cw = Inches(2.9); gx = Inches(0.2); x0 = Inches(0.7); y0 = Inches(1.35)
for i, (fname, name, role, blurb) in enumerate(members):
    l = x0 + i * (cw + gx)
    picture_fit(s, img(fname), l, y0, cw, Inches(2.6))
    textbox(s, l, y0 + Inches(2.7), cw, Inches(0.4), name, size=15, bold=True,
            color=DARK, align=PP_ALIGN.CENTER)
    textbox(s, l, y0 + Inches(3.1), cw, Inches(0.35), role, size=12.5,
            color=RED, align=PP_ALIGN.CENTER, italic=True)
    textbox(s, l, y0 + Inches(3.5), cw, Inches(1.6), blurb, size=11.5,
            color=GREY, align=PP_ALIGN.CENTER)
textbox(s, Inches(0.7), Inches(6.9), Inches(11.9), Inches(0.4),
    "Team The Red Castle  —  HMK AI and Robotics Club", size=13, bold=True,
    color=DARK, align=PP_ALIGN.CENTER)

# ==========================================================================
# 8. LINKS AND CREDITS
# ==========================================================================
s = add_slide(); bg(s)
header_bar(s, "8. Links and Credits", number=num())
bullets(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(2.8), [
    "GitHub repository (full source, CAD, documentation):",
    "  https://github.com/jolianjij/Red-Castle-WRO-Future-Engineers-2026",
    "Open Challenge driving video:",
    "  https://youtu.be/GmrX7NhcAxE",
    "Obstacle Challenge driving video:",
    "  https://youtu.be/ujbImlpQH6g",
], size=16, gap=10)
textbox(s, Inches(0.7), Inches(4.6), Inches(11.9), Inches(0.4),
    "Thanks", size=17, bold=True, color=RED)
textbox(s, Inches(0.7), Inches(5.05), Inches(11.9), Inches(1.8),
    "As a team, we thank HMK AI and Robotics Club for hosting and supporting us "
    "throughout this build, and our coach, Ahmad Kalthom, for the guidance that "
    "kept the project moving through every dead end along the way.",
    size=15, color=DARK)

prs.save(os.path.join(ROOT, "Engineering-Journal.pptx"))
print("saved -> Engineering-Journal.pptx  (%d slides)" % len(prs.slides))
